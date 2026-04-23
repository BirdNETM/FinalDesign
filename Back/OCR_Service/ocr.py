import os

# 🔴 必须在 import paddleocr 之前
os.environ["FLAGS_use_mkldnn"] = "0"
os.environ["FLAGS_enable_pir_api"] = "0"
os.environ["FLAGS_new_executor"] = "0"


from fastapi import FastAPI, Query, HTTPException, Body
import docx
# 读取doc文件需要额外库（pip install pywin32，仅Windows可用）
try:
    import win32com.client
except ImportError:
    win32com = None
from pathlib import Path
import numpy as np
from paddleocr import PaddleOCR
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from io import BytesIO

from db import get_cursor, commit

ocr = PaddleOCR(use_angle_cls=True, lang='ch', enable_mkldnn=False)  # 中文

def ocr_image(img: Image.Image):
    """将 PIL.Image 转为 numpy array 后用 PaddleOCR 识别文字"""
    img_rgb = img.convert("RGB")  # 确保是 RGB
    img_np = np.array(img_rgb)
    try:
        result = ocr.ocr(img_np)
    except Exception as e:
        return [f"OCR 识别失败: {e}"]

    texts = []
    for line in result:
        for box in line:
            texts.append(box[1][0])
    return texts

def ocr_file(doc_id):
    cursor = get_cursor()
    sql = "select doc_url from docs where doc_id = %s"
    cursor.execute(sql, (
        doc_id,
    ))
    result = cursor.fetchone()
    path = Path(result[0])
    print(path)
    commit()

    if not path.exists():
        raise HTTPException(status_code=400, detail="File not found")

    suffix = path.suffix.lower().lstrip(".")
    texts = []
    if suffix in ["doc", "docx", "txt"]:
        try:
            # 读取txt文件
            if suffix == "txt":
                # 兼容常见编码：utf-8、gbk（Windows中文文本）
                encodings = ["utf-8", "gbk", "gb2312"]
                for enc in encodings:
                    try:
                        with open(path, "r", encoding=enc) as f:
                            texts = [line.strip() for line in f if line.strip()]
                        break
                    except UnicodeDecodeError:
                        continue

            # 读取docx文件
            elif suffix == "docx":
                doc = docx.Document(path)
                # 提取所有段落文本
                texts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]

            # 读取doc文件（仅Windows系统可用）
            elif suffix == "doc":
                if not win32com:
                    raise HTTPException(status_code=500, detail="请安装pywin32库：pip install pywin32")
                # 调用Word程序读取doc
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(str(path))
                texts = [doc.Content.Text.strip()]
                doc.Close()
                word.Quit()

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"读取文件失败: {str(e)}")

    elif suffix in ["jpg", "jpeg", "png"]:
        try:
            img = Image.open(path)
            texts += ocr_image(img)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"图片打开或识别失败: {e}")

    elif suffix == "pdf":
        try:
            pages = convert_from_path(
                str(path),
                dpi=300,
                poppler_path=r"D:\complier\popular\poppler-25.12.0\Library\bin"
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PDF 转图片失败: {e}")

        for idx, page in enumerate(pages, start=1):
            page_texts = ocr_image(page)
            # texts.append(f"=== 第 {idx} 页 ===")
            texts.extend(page_texts)

    elif suffix in ["ppt", "pptx"]:
        try:
            prs = Presentation(path)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPT 打开失败: {e}")

        for idx, slide in enumerate(prs.slides, start=1):
            # texts.append(f"=== 第 {idx} 张幻灯片 ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
                elif getattr(shape, "image", None) is not None:
                    try:
                        img_bytes = shape.image.blob
                        img = Image.open(BytesIO(img_bytes))
                        texts.extend(ocr_image(img))
                    except Exception:
                        pass
    else:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式: {suffix}")
    return "\n".join(texts)